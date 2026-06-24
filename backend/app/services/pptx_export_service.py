"""Service for exporting slides artifacts to PowerPoint (PPTX) format."""

from io import BytesIO
from typing import Any, Dict, List

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR


class PptxExportService:
    """Service for generating PowerPoint presentations from slides data."""

    # Slide dimensions (widescreen 16:9)
    SLIDE_WIDTH = Inches(13.333)
    SLIDE_HEIGHT = Inches(7.5)

    # Color scheme
    COLORS = {
        "dark_bg": RGBColor(15, 23, 42),        # slate-900
        "light_text": RGBColor(248, 250, 252),  # slate-50
        "accent": RGBColor(59, 130, 246),       # blue-500
        "muted": RGBColor(148, 163, 184),       # slate-400
        "success": RGBColor(34, 197, 94),       # green-500
        "danger": RGBColor(239, 68, 68),        # red-500
    }

    def generate_pptx(self, slides: List[Dict[str, Any]], title: str = "Presentation") -> BytesIO:
        """Generate a PPTX file from slides data.

        Args:
            slides: List of slide data dictionaries
            title: Presentation title

        Returns:
            BytesIO buffer containing the PPTX file
        """
        prs = Presentation()
        prs.slide_width = self.SLIDE_WIDTH
        prs.slide_height = self.SLIDE_HEIGHT

        for slide_data in slides:
            slide_type = slide_data.get("type", "text")

            if slide_type == "title":
                self._add_title_slide(prs, slide_data)
            elif slide_type == "metrics":
                self._add_metrics_slide(prs, slide_data)
            elif slide_type == "chart":
                self._add_chart_slide(prs, slide_data)
            elif slide_type == "bullets":
                self._add_bullets_slide(prs, slide_data)
            elif slide_type == "code":
                self._add_code_slide(prs, slide_data)
            elif slide_type == "text":
                self._add_text_slide(prs, slide_data)
            else:
                # Default to text slide for unknown types
                self._add_text_slide(prs, slide_data)

        # Save to buffer
        buffer = BytesIO()
        prs.save(buffer)
        buffer.seek(0)
        return buffer

    def _add_title_slide(self, prs: Presentation, data: Dict[str, Any]) -> None:
        """Add a title slide."""
        blank_layout = prs.slide_layouts[6]  # Blank layout
        slide = prs.slides.add_slide(blank_layout)

        # Dark background
        background = slide.shapes.add_shape(
            1,  # Rectangle
            Inches(0), Inches(0),
            self.SLIDE_WIDTH, self.SLIDE_HEIGHT
        )
        background.fill.solid()
        background.fill.fore_color.rgb = self.COLORS["dark_bg"]
        background.line.fill.background()

        # Title - allow for wrapping with larger box
        title_text = data.get("title", "Untitled")
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(2),
            Inches(12.333), Inches(2)
        )
        title_frame = title_box.text_frame
        title_frame.word_wrap = True
        title_para = title_frame.paragraphs[0]
        title_para.text = title_text
        # Adjust font size based on title length
        if len(title_text) > 50:
            title_para.font.size = Pt(40)
        else:
            title_para.font.size = Pt(54)
        title_para.font.bold = True
        title_para.font.color.rgb = self.COLORS["light_text"]
        title_para.alignment = PP_ALIGN.CENTER

        # Subtitle - positioned below title with more spacing
        if data.get("subtitle"):
            subtitle_box = slide.shapes.add_textbox(
                Inches(0.5), Inches(4.5),
                Inches(12.333), Inches(1)
            )
            subtitle_frame = subtitle_box.text_frame
            subtitle_frame.word_wrap = True
            subtitle_para = subtitle_frame.paragraphs[0]
            subtitle_para.text = data.get("subtitle", "")
            subtitle_para.font.size = Pt(28)
            subtitle_para.font.color.rgb = self.COLORS["muted"]
            subtitle_para.alignment = PP_ALIGN.CENTER

    def _add_metrics_slide(self, prs: Presentation, data: Dict[str, Any]) -> None:
        """Add a metrics/KPI slide."""
        blank_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_layout)

        # Dark background
        background = slide.shapes.add_shape(
            1, Inches(0), Inches(0),
            self.SLIDE_WIDTH, self.SLIDE_HEIGHT
        )
        background.fill.solid()
        background.fill.fore_color.rgb = self.COLORS["dark_bg"]
        background.line.fill.background()

        # Slide title
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.5),
            Inches(12.333), Inches(1)
        )
        title_frame = title_box.text_frame
        title_para = title_frame.paragraphs[0]
        title_para.text = data.get("title", "Key Metrics")
        title_para.font.size = Pt(36)
        title_para.font.bold = True
        title_para.font.color.rgb = self.COLORS["light_text"]

        # Metrics cards
        metrics = data.get("metrics", [])
        if not metrics:
            return

        # Calculate card layout
        num_metrics = min(len(metrics), 4)  # Max 4 metrics per slide
        card_width = Inches(2.8)
        card_height = Inches(2)
        total_width = num_metrics * card_width + (num_metrics - 1) * Inches(0.3)
        start_x = (self.SLIDE_WIDTH - total_width) / 2

        for i, metric in enumerate(metrics[:4]):
            x = start_x + i * (card_width + Inches(0.3))
            y = Inches(3)

            # Metric value
            value_box = slide.shapes.add_textbox(x, y, card_width, Inches(1))
            value_frame = value_box.text_frame
            value_para = value_frame.paragraphs[0]
            value_para.text = str(metric.get("value", "N/A"))
            value_para.font.size = Pt(44)
            value_para.font.bold = True
            value_para.font.color.rgb = self.COLORS["light_text"]
            value_para.alignment = PP_ALIGN.CENTER

            # Metric label
            label_box = slide.shapes.add_textbox(x, y + Inches(1), card_width, Inches(0.5))
            label_frame = label_box.text_frame
            label_para = label_frame.paragraphs[0]
            label_para.text = metric.get("label", "")
            label_para.font.size = Pt(18)
            label_para.font.color.rgb = self.COLORS["muted"]
            label_para.alignment = PP_ALIGN.CENTER

            # Change indicator
            change = metric.get("change")
            if change:
                change_box = slide.shapes.add_textbox(x, y + Inches(1.5), card_width, Inches(0.4))
                change_frame = change_box.text_frame
                change_para = change_frame.paragraphs[0]
                change_para.text = change
                change_para.font.size = Pt(16)
                change_para.font.bold = True
                # Color based on positive/negative
                if change.startswith("+") or "increase" in change.lower():
                    change_para.font.color.rgb = self.COLORS["success"]
                elif change.startswith("-") or "decrease" in change.lower():
                    change_para.font.color.rgb = self.COLORS["danger"]
                else:
                    change_para.font.color.rgb = self.COLORS["accent"]
                change_para.alignment = PP_ALIGN.CENTER

    def _add_chart_slide(self, prs: Presentation, data: Dict[str, Any]) -> None:
        """Add a chart slide (placeholder - charts need to be added manually or via chart API)."""
        blank_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_layout)

        # Dark background
        background = slide.shapes.add_shape(
            1, Inches(0), Inches(0),
            self.SLIDE_WIDTH, self.SLIDE_HEIGHT
        )
        background.fill.solid()
        background.fill.fore_color.rgb = self.COLORS["dark_bg"]
        background.line.fill.background()

        # Slide title
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.5),
            Inches(12.333), Inches(1)
        )
        title_frame = title_box.text_frame
        title_para = title_frame.paragraphs[0]
        title_para.text = data.get("title", "Chart")
        title_para.font.size = Pt(36)
        title_para.font.bold = True
        title_para.font.color.rgb = self.COLORS["light_text"]

        # Chart placeholder
        chart_placeholder = slide.shapes.add_shape(
            1,  # Rectangle
            Inches(1), Inches(2),
            Inches(11.333), Inches(4.5)
        )
        chart_placeholder.fill.solid()
        chart_placeholder.fill.fore_color.rgb = RGBColor(30, 41, 59)  # slate-800
        chart_placeholder.line.color.rgb = RGBColor(51, 65, 85)  # slate-700

        # Chart type label
        chart_type = data.get("chartType", "chart")
        label_box = slide.shapes.add_textbox(
            Inches(1), Inches(3.75),
            Inches(11.333), Inches(1)
        )
        label_frame = label_box.text_frame
        label_para = label_frame.paragraphs[0]
        label_para.text = f"[{chart_type.upper()} CHART]"
        label_para.font.size = Pt(24)
        label_para.font.color.rgb = self.COLORS["muted"]
        label_para.alignment = PP_ALIGN.CENTER

        # Insight text
        insight = data.get("insight")
        if insight:
            insight_box = slide.shapes.add_textbox(
                Inches(1), Inches(6.7),
                Inches(11.333), Inches(0.6)
            )
            insight_frame = insight_box.text_frame
            insight_para = insight_frame.paragraphs[0]
            insight_para.text = insight
            insight_para.font.size = Pt(18)
            insight_para.font.italic = True
            insight_para.font.color.rgb = self.COLORS["accent"]
            insight_para.alignment = PP_ALIGN.CENTER

    def _add_bullets_slide(self, prs: Presentation, data: Dict[str, Any]) -> None:
        """Add a bullet points slide."""
        blank_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_layout)

        # Dark background
        background = slide.shapes.add_shape(
            1, Inches(0), Inches(0),
            self.SLIDE_WIDTH, self.SLIDE_HEIGHT
        )
        background.fill.solid()
        background.fill.fore_color.rgb = self.COLORS["dark_bg"]
        background.line.fill.background()

        # Slide title
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.5),
            Inches(12.333), Inches(1)
        )
        title_frame = title_box.text_frame
        title_para = title_frame.paragraphs[0]
        title_para.text = data.get("title", "Key Points")
        title_para.font.size = Pt(36)
        title_para.font.bold = True
        title_para.font.color.rgb = self.COLORS["light_text"]

        # Bullet points
        bullets = data.get("bullets", [])
        if not bullets:
            return

        bullets_box = slide.shapes.add_textbox(
            Inches(1), Inches(2),
            Inches(11.333), Inches(5)
        )
        bullets_frame = bullets_box.text_frame
        bullets_frame.word_wrap = True

        for i, bullet in enumerate(bullets[:8]):  # Max 8 bullets
            if i == 0:
                para = bullets_frame.paragraphs[0]
            else:
                para = bullets_frame.add_paragraph()

            para.text = f"• {bullet}"
            para.font.size = Pt(24)
            para.font.color.rgb = self.COLORS["light_text"]
            para.space_before = Pt(12)
            para.space_after = Pt(12)

        # Insight text if present
        insight = data.get("insight")
        if insight:
            insight_box = slide.shapes.add_textbox(
                Inches(1), Inches(6.5),
                Inches(11.333), Inches(0.6)
            )
            insight_frame = insight_box.text_frame
            insight_para = insight_frame.paragraphs[0]
            insight_para.text = insight
            insight_para.font.size = Pt(18)
            insight_para.font.italic = True
            insight_para.font.color.rgb = self.COLORS["accent"]
            insight_para.alignment = PP_ALIGN.CENTER

    def _add_code_slide(self, prs: Presentation, data: Dict[str, Any]) -> None:
        """Add a code/methodology slide with code snippets."""
        blank_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_layout)

        # Dark background
        background = slide.shapes.add_shape(
            1, Inches(0), Inches(0),
            self.SLIDE_WIDTH, self.SLIDE_HEIGHT
        )
        background.fill.solid()
        background.fill.fore_color.rgb = self.COLORS["dark_bg"]
        background.line.fill.background()

        # Title
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.5),
            Inches(12.333), Inches(1)
        )
        title_frame = title_box.text_frame
        title_para = title_frame.paragraphs[0]
        title_para.text = data.get("title", "Methodology")
        title_para.font.size = Pt(36)
        title_para.font.bold = True
        title_para.font.color.rgb = self.COLORS["light_text"]

        # Code snippets
        code_snippets = data.get("code_snippets", [])
        y_pos = Inches(1.8)

        for snippet in code_snippets[:3]:  # Max 3 snippets
            # Code box background
            code_bg = slide.shapes.add_shape(
                1,  # Rectangle
                Inches(0.5), y_pos,
                Inches(12.333), Inches(1.4)
            )
            code_bg.fill.solid()
            code_bg.fill.fore_color.rgb = RGBColor(30, 41, 59)  # slate-800
            code_bg.line.color.rgb = RGBColor(51, 65, 85)  # slate-700

            # Code text
            code_box = slide.shapes.add_textbox(
                Inches(0.7), y_pos + Inches(0.15),
                Inches(11.933), Inches(1.2)
            )
            code_frame = code_box.text_frame
            code_frame.word_wrap = True
            code_para = code_frame.paragraphs[0]
            code_para.text = snippet[:200]  # Truncate long snippets
            code_para.font.size = Pt(14)
            code_para.font.name = "Consolas"
            code_para.font.color.rgb = self.COLORS["light_text"]

            y_pos += Inches(1.6)

        # Additional text if present
        text_content = data.get("text")
        if text_content and y_pos < Inches(6):
            text_box = slide.shapes.add_textbox(
                Inches(0.5), y_pos + Inches(0.2),
                Inches(12.333), Inches(1.5)
            )
            text_frame = text_box.text_frame
            text_frame.word_wrap = True
            text_para = text_frame.paragraphs[0]
            text_para.text = text_content[:300]
            text_para.font.size = Pt(18)
            text_para.font.color.rgb = self.COLORS["muted"]

    def _add_text_slide(self, prs: Presentation, data: Dict[str, Any]) -> None:
        """Add a generic text slide."""
        blank_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_layout)

        # Dark background
        background = slide.shapes.add_shape(
            1, Inches(0), Inches(0),
            self.SLIDE_WIDTH, self.SLIDE_HEIGHT
        )
        background.fill.solid()
        background.fill.fore_color.rgb = self.COLORS["dark_bg"]
        background.line.fill.background()

        # Title
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.5),
            Inches(12.333), Inches(1)
        )
        title_frame = title_box.text_frame
        title_para = title_frame.paragraphs[0]
        title_para.text = data.get("title", "")
        title_para.font.size = Pt(36)
        title_para.font.bold = True
        title_para.font.color.rgb = self.COLORS["light_text"]

        # Content text
        content = data.get("content") or data.get("text") or ""
        if content:
            content_box = slide.shapes.add_textbox(
                Inches(0.5), Inches(2),
                Inches(12.333), Inches(4)
            )
            content_frame = content_box.text_frame
            content_frame.word_wrap = True
            content_para = content_frame.paragraphs[0]
            content_para.text = content
            content_para.font.size = Pt(22)
            content_para.font.color.rgb = self.COLORS["light_text"]

        # Insight text if present
        insight = data.get("insight")
        if insight:
            insight_box = slide.shapes.add_textbox(
                Inches(0.5), Inches(6.5),
                Inches(12.333), Inches(0.6)
            )
            insight_frame = insight_box.text_frame
            insight_para = insight_frame.paragraphs[0]
            insight_para.text = insight
            insight_para.font.size = Pt(18)
            insight_para.font.italic = True
            insight_para.font.color.rgb = self.COLORS["accent"]
            insight_para.alignment = PP_ALIGN.CENTER
