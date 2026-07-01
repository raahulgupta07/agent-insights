"""Pipeline v1 (P2, HYBRID_LOGIC_PARSER): parse a logic / Q&A document into
structured triples the training pipeline can verify against.

A logic doc (e.g. ``CRM Agent Q&A , Logic.docx``) states, per business question:

    Q1.  What is the total number of leads across all channels?
    Ans: 1,544.
    Logic: Date = Call Completed Date, Status = Completed,
           Call Outcome = Unsuccessful, Related Brand Relationship: Type = Lead

This module turns that into::

    {"question": "...", "answer_text": "1,544.", "expected": 1544,
     "logic_text": "Status = Completed; Call Outcome = Unsuccessful; Type = Lead",
     "filters": [{"column": "Status", "op": "=", "value": "Completed"}, ...]}

The triples feed the Definition Registry (P3) + eval gate (P5): ``expected`` is
the ground-truth number a generated query must reproduce, ``filters`` are the
documented predicate. Pure + fail-soft — never raises into the caller; returns
[] on any problem.
"""
from __future__ import annotations

import logging
import re
from typing import Dict, List, Optional

logger = logging.getLogger(__name__)

# A line that starts a question: "Q1.", "Q 4.", "Q12)" ...
_Q_RE = re.compile(r"^\s*Q\s*\.?\s*(\d+)\s*[\.\):]", re.IGNORECASE)
_ANS_RE = re.compile(r"^\s*ans\w*\s*[:\-]", re.IGNORECASE)
_LOGIC_RE = re.compile(r"^\s*logic\s*[:\-]", re.IGNORECASE)
# an integer with thousands separators, OR a bare integer
_INT_RE = re.compile(r"\d{1,3}(?:[,\s]\d{3})+|\d+")
# first integer (with thousands separators) inside an answer line
_NUM_RE = re.compile(r"(\d{1,3}(?:[,\s]\d{3})+|\d+)")

# ratio blocks: "Numerator(s) : ( ... )" / "Denominator(s): ( ... )" (P8 RATIO_METRICS).
# The keyword may be singular/plural, followed by ':' (with optional space) then a
# parenthesised predicate list. Column names in the doc themselves contain ':', so the
# split of each predicate is left to _filters_from_logic (splits on the FIRST '=' only).
_NUM_BLOCK_RE = re.compile(
    r"numerators?\s*:?\s*\(([^)]*)\)", re.IGNORECASE)
_DEN_BLOCK_RE = re.compile(
    r"denominators?\s*:?\s*\(([^)]*)\)", re.IGNORECASE)


# --------------------------------------------------------------------------- #
# text extraction per file type
# --------------------------------------------------------------------------- #

def extract_text(path: str) -> str:
    """Return the document's plain text (docx/pptx/pdf/txt). '' on failure."""
    ext = (path.rsplit(".", 1)[-1] if "." in path else "").lower()
    try:
        if ext in ("docx",):
            import docx

            d = docx.Document(path)
            lines = [p.text for p in d.paragraphs]
            for tbl in getattr(d, "tables", []):
                for row in tbl.rows:
                    lines.append("\t".join(c.text for c in row.cells))
            return "\n".join(lines)
        if ext in ("pptx",):
            from pptx import Presentation

            prs = Presentation(path)
            out = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        out.append(shape.text_frame.text)
            return "\n".join(out)
        if ext in ("pdf",):
            from pypdf import PdfReader

            r = PdfReader(path)
            return "\n".join((pg.extract_text() or "") for pg in r.pages)
        # txt / fallback
        with open(path, "r", encoding="utf-8", errors="replace") as f:
            return f.read()
    except Exception:  # noqa: BLE001
        logger.warning("logic_parser.extract_text failed for %s", path, exc_info=True)
        return ""


# --------------------------------------------------------------------------- #
# triple parsing
# --------------------------------------------------------------------------- #

def _int_val(tok: str) -> Optional[int]:
    try:
        return int(re.sub(r"[,\s]", "", tok))
    except Exception:  # noqa: BLE001
        return None


def _is_bare_year(tok: str, n: int) -> bool:
    """A bare 4-digit token in 1900–2100 is almost certainly a year, not a count."""
    return len(re.sub(r"[,\s]", "", tok)) == 4 and 1900 <= n <= 2100


def _expected_from_answer(answer: str) -> Optional[int]:
    """The ground-truth count in the answer line.

    Prefers the number written after 'total' or '=' (that's the metric value);
    otherwise the first integer — but skips a bare 4-digit year (1900–2100) like
    the '2025' in 'Jan–Jun 2025 total 644' so we don't return the year as the count.
    Handles '1,544' / '7 526'. None if no usable number.
    """
    try:
        s = answer or ""
        # 1) a number right after 'total' or '=' wins ('... total 644', 'x = 644')
        m = re.search(r"(?:total|=)\D{0,12}?(" + _INT_RE.pattern + r")", s, re.IGNORECASE)
        if m:
            v = _int_val(m.group(1))
            if v is not None:
                return v
        # 2) otherwise first integer that isn't a bare year
        for tok in _INT_RE.findall(s):
            v = _int_val(tok)
            if v is None:
                continue
            if _is_bare_year(tok, v):
                continue
            return v
        return None
    except Exception:  # noqa: BLE001
        return None


def _filters_from_logic(logic: str) -> List[Dict[str, str]]:
    """Split a logic line into {column, op, value} predicates.

    Handles 'Status = Completed', 'Call Outcome=Unsuccessful',
    'Related Brand Relationship: Type = Lead' (column may itself contain ':').
    Conservative — skips the bare 'Date = Call Completed Date' grouping hint.
    """
    out: List[Dict[str, str]] = []
    if not logic:
        return out
    # split on commas that separate predicates
    for part in re.split(r"\s*,\s*", logic):
        part = part.strip().rstrip(".")
        if "=" not in part:
            # a comma-tail with no '=' (e.g. 'Existing' in 'Status = Retained, Existing')
            # is another value for the previous equality predicate -> OR it in with '/'
            # so build_predicate expands to an IN (...) list.
            tail = part.strip().strip("'\"").strip(" .;").strip()
            if tail and out and out[-1].get("op") == "=":
                out[-1]["value"] = out[-1]["value"] + " / " + tail
            continue
        col, val = part.split("=", 1)
        col = col.strip()
        # strip surrounding quotes + trailing sentence punctuation/space so a value
        # like "New." (end of a sentence) becomes "New".
        val = val.strip().strip("'\"").strip(" .;").strip()
        if not col or not val:
            continue
        low = col.lower()
        # 'Date = Call Completed Date' is a grouping hint, not an equality filter
        if low in ("date", "group by", "groupby") or val.lower().endswith("date"):
            out.append({"column": col, "op": "groupby", "value": val})
            continue
        out.append({"column": col, "op": "=", "value": val})
    return out


def _ratio_from_logic(text: str) -> Optional[Dict]:
    """If the logic/answer text carries Numerator(s)+Denominator(s) blocks, return
    ``{"num_filters": [...], "den_filters": [...], "group_by": [...]}``; else None.

    Each block's parenthesised body is split with the same predicate splitter used
    for a flat metric, so column names that themselves contain ':' are preserved and
    'A / B' or 'A, B' value sets expand downstream. Group-by dims are any 'groupby'
    hints found OUTSIDE the two blocks (e.g. 'Date = Call Completed Date'). Fail-soft.
    """
    try:
        if not text:
            return None
        num_m = _NUM_BLOCK_RE.search(text)
        den_m = _DEN_BLOCK_RE.search(text)
        if not num_m or not den_m:
            return None
        num_filters = _filters_from_logic(num_m.group(1))
        den_filters = _filters_from_logic(den_m.group(1))
        if not num_filters or not den_filters:
            return None
        # group-by dims = groupby hints in the logic outside the num/den parens
        outside = _NUM_BLOCK_RE.sub("", _DEN_BLOCK_RE.sub("", text))
        group_by = [f for f in _filters_from_logic(outside) if f.get("op") == "groupby"]
        return {"num_filters": num_filters, "den_filters": den_filters,
                "group_by": group_by}
    except Exception:  # noqa: BLE001
        logger.warning("logic_parser._ratio_from_logic failed", exc_info=True)
        return None


def parse_triples(text: str) -> List[Dict]:
    """Parse Q/Ans/Logic blocks out of the document text. Never raises."""
    triples: List[Dict] = []
    try:
        lines = [ln.rstrip() for ln in (text or "").splitlines()]
        cur: Optional[Dict] = None
        mode = None  # which field we're currently appending continuation lines to
        for ln in lines:
            if _Q_RE.match(ln):
                if cur and cur.get("question"):
                    triples.append(cur)
                qn = _Q_RE.match(ln).group(1)
                cur = {"n": int(qn), "question": _Q_RE.sub("", ln).strip(),
                       "answer_text": "", "logic_text": ""}
                mode = "question"
            elif cur is not None and _ANS_RE.match(ln):
                cur["answer_text"] = _ANS_RE.sub("", ln).strip()
                mode = "answer"
            elif cur is not None and _LOGIC_RE.match(ln):
                cur["logic_text"] = _LOGIC_RE.sub("", ln).strip()
                mode = "logic"
            elif cur is not None and ln.strip() and mode in ("logic", "answer"):
                # continuation of a wrapped logic/answer line
                cur[f"{mode}_text"] = (cur.get(f"{mode}_text", "") + " " + ln.strip()).strip()
        if cur and cur.get("question"):
            triples.append(cur)

        # enrich with expected number + parsed filters
        ratio_on = False
        try:
            from app.settings.hybrid_flags import flags
            ratio_on = bool(flags.RATIO_METRICS)
        except Exception:  # noqa: BLE001
            ratio_on = False
        for t in triples:
            t["expected"] = _expected_from_answer(t.get("answer_text", ""))
            t["filters"] = _filters_from_logic(t.get("logic_text", ""))
            # ratio metrics (P8): a Numerator+Denominator pair -> tag + split blocks.
            # Scan logic and answer text (Q11 states num/den in the answer). Flag-gated;
            # off -> untouched flat metric, so old behaviour is byte-identical.
            if ratio_on:
                r = _ratio_from_logic(
                    (t.get("logic_text", "") + " " + t.get("answer_text", "")).strip())
                if r:
                    t["metric_kind"] = "ratio"
                    t["num_filters"] = r["num_filters"]
                    t["den_filters"] = r["den_filters"]
                    t["group_by"] = r["group_by"]
    except Exception:  # noqa: BLE001
        logger.warning("logic_parser.parse_triples failed", exc_info=True)
    return triples


def parse_logic_doc(path: str) -> List[Dict]:
    """Top-level: extract text from the file then parse Q/Ans/Logic triples.

    Returns a list of {n, question, answer_text, expected, logic_text, filters}.
    Fail-soft -> [].
    """
    text = extract_text(path)
    if not text:
        return []
    triples = parse_triples(text)
    logger.info("logic_parser: %s -> %d triple(s), %d with expected number",
                path.rsplit("/", 1)[-1], len(triples),
                sum(1 for t in triples if t.get("expected") is not None))
    return triples


def looks_like_logic_doc(text: str) -> bool:
    """Heuristic: text has >=2 Q-markers AND a Logic: or Ans: marker."""
    try:
        qn = len(re.findall(r"(?im)^\s*Q\s*\.?\s*\d+\s*[\.\):]", text or ""))
        has_logic = bool(re.search(r"(?im)^\s*logic\s*[:\-]", text or ""))
        has_ans = bool(re.search(r"(?im)^\s*ans\w*\s*[:\-]", text or ""))
        return qn >= 2 and (has_logic or has_ans)
    except Exception:  # noqa: BLE001
        return False
