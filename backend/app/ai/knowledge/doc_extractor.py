"""Auto-configure-from-doc extractor (AUTOMAP).

Parses an uploaded *definitions spreadsheet* (.xlsx) and/or an *explanation
deck* (.pptx), builds a single text digest, asks the configured LLM (OpenRouter)
to extract a structured proposal, and fuzzy-matches the proposed column names
against a target data source's live schema (its ``DataSourceTable`` rows).

The proposal shape returned to the route::

    {
        "column_descriptions": [
            {"column": "...", "description": "...",
             "matched": bool, "table_id": str|None, "table_name": str|None,
             "matched_column": str|None, "match_kind": "exact|ci|strip|fuzzy|none"}
        ],
        "instructions":  [{"content": "...", "category": "..."}],
        "examples":      [{"question": "...", "answer": "...", "sql": "..."}],
        "compliance":    [{"content": "..."}],
        "unmatched_columns": ["..."],   # proposed columns with no schema match
        "schema_columns":    ["table.col", ...],  # what the live schema offered
        "source": "llm" | "fallback",
        "warnings": ["..."],
    }

Design rules (mirror the rest of the hybrid layer):
- Never raise into the route on *content* problems — return a dict with an
  ``error`` key so the route can answer a helpful 422. Hard programmer errors
  (bad args) may still raise.
- The parser libs (``openpyxl`` / ``python-pptx``) may NOT be installed in the
  runtime container. Imports are wrapped; a missing lib yields a clear
  ``{"error": "openpyxl not installed"}`` style result.
- LLM is OpenRouter-only via the repo idiom
  ``LLM(model, usage_session_maker=...).inference(prompt, usage_scope=...)``.
- Deterministic OFFLINE fallback: if no LLM/model is available, parse a 2-column
  definitions sheet directly into ``column_descriptions`` so the feature still
  works without a model.
"""

import difflib
import json
import logging
import os
from typing import Any, Dict, List, Optional, Tuple

from sqlalchemy import select
from sqlalchemy.ext.asyncio import AsyncSession

from app.models.data_source import DataSource
from app.models.datasource_table import DataSourceTable
from app.models.file import File as FileModel

logger = logging.getLogger(__name__)

# Cap how much digest text we send to the model (cheap + bounded).
_MAX_DIGEST_CHARS = 16000
# difflib fuzzy-match cutoff.
_FUZZY_CUTOFF = 0.6


# --------------------------------------------------------------------------- #
# File-path resolution (mirrors data_source_from_file._resolve_upload_path)
# --------------------------------------------------------------------------- #
def _resolve_upload_path(stored_path: str) -> Optional[str]:
    """Resolve the on-disk absolute path for an uploaded file, traversal-safe.

    Uploaded files live flat under ``<cwd>/uploads/files/<basename>`` (see
    routes/file.py). In-container that is ``/app/backend/uploads/files/``.
    Returns the path if it exists, else None.
    """
    if not stored_path:
        return None
    base = os.path.basename(stored_path)
    candidate = os.path.join(os.getcwd(), "uploads", "files", base)
    if os.path.exists(candidate):
        return candidate
    rel = os.path.join(os.getcwd(), stored_path)
    if os.path.exists(rel):
        return rel
    if os.path.isabs(stored_path) and os.path.exists(stored_path):
        return stored_path
    return None


# --------------------------------------------------------------------------- #
# Parsers (libs may be missing — fail soft with a clear error dict)
# --------------------------------------------------------------------------- #
def _parse_xlsx(path: str) -> Dict[str, Any]:
    """Parse an .xlsx into {digest, rows} or {error}.

    ``rows`` is a list of (left, right) 2-col pairs harvested from every sheet
    (used by the offline fallback as column -> definition). ``digest`` is a
    readable text rendering of the workbook for the LLM.
    """
    try:
        import openpyxl  # type: ignore
    except Exception:
        return {"error": "openpyxl not installed"}

    try:
        wb = openpyxl.load_workbook(path, read_only=True, data_only=True)
    except Exception as e:
        return {"error": f"could not read spreadsheet: {e}"}

    lines: List[str] = []
    pairs: List[Tuple[str, str]] = []
    try:
        for ws in wb.worksheets:
            lines.append(f"# Sheet: {ws.title}")
            for row in ws.iter_rows(values_only=True):
                cells = [
                    str(c).strip() for c in row
                    if c is not None and str(c).strip() != ""
                ]
                if not cells:
                    continue
                lines.append(" | ".join(cells))
                # 2-col definition pair (col name -> definition).
                if len(cells) >= 2 and cells[0]:
                    pairs.append((cells[0], cells[1]))
    finally:
        try:
            wb.close()
        except Exception:
            pass

    return {"digest": "\n".join(lines), "pairs": pairs}


def _parse_pptx(path: str) -> Dict[str, Any]:
    """Parse a .pptx into {digest} or {error}. All slide text, in order."""
    try:
        from pptx import Presentation  # type: ignore
    except Exception:
        return {"error": "python-pptx not installed"}

    try:
        prs = Presentation(path)
    except Exception as e:
        return {"error": f"could not read presentation: {e}"}

    lines: List[str] = []
    for idx, slide in enumerate(prs.slides, start=1):
        lines.append(f"# Slide {idx}")
        for shape in slide.shapes:
            text = getattr(shape, "text", None)
            if text and text.strip():
                lines.append(text.strip())
            # Tables on a slide (column->definition decks).
            if getattr(shape, "has_table", False):
                try:
                    for trow in shape.table.rows:
                        cells = [c.text.strip() for c in trow.cells if c.text]
                        if cells:
                            lines.append(" | ".join(cells))
                except Exception:
                    pass
    return {"digest": "\n".join(lines)}


def _parse_file(path: str) -> Dict[str, Any]:
    """Dispatch on extension. Returns {digest,pairs?} or {error}."""
    ext = os.path.splitext(path)[1].lower()
    if ext in (".xlsx", ".xlsm"):
        return _parse_xlsx(path)
    if ext in (".pptx",):
        return _parse_pptx(path)
    return {"error": f"unsupported file type '{ext}' (need .xlsx or .pptx)"}


# --------------------------------------------------------------------------- #
# Schema loading + fuzzy column matching
# --------------------------------------------------------------------------- #
async def _load_schema_columns(
    db: AsyncSession, data_source_id: str
) -> List[Dict[str, Any]]:
    """Load the live columns of a data source's active tables.

    Returns a flat list of
    ``{"table_id", "table_name", "column", "dtype"}`` from each
    ``DataSourceTable.columns`` JSON (list of ``{name, dtype, description?}``).
    Falls back to the ConnectionTable schema columns when the legacy
    ``columns`` JSON is empty.
    """
    res = await db.execute(
        select(DataSourceTable).where(
            DataSourceTable.datasource_id == str(data_source_id),
            DataSourceTable.is_active == True,  # noqa: E712
        )
    )
    out: List[Dict[str, Any]] = []
    for t in res.scalars().all():
        cols = t.columns or []
        if not cols:
            # Fall back to the prompt-table view (ConnectionTable-backed).
            try:
                pt = t.to_prompt_table()
                cols = [
                    {"name": c.name, "dtype": getattr(c, "dtype", None)}
                    for c in (pt.columns or [])
                ]
            except Exception:
                cols = []
        for c in cols:
            if not isinstance(c, dict):
                continue
            name = c.get("name")
            if not name:
                continue
            out.append({
                "table_id": str(t.id),
                "table_name": t.name,
                "column": name,
                "dtype": c.get("dtype"),
            })
    return out


def _match_column(
    proposed: str, schema_cols: List[Dict[str, Any]]
) -> Tuple[Optional[Dict[str, Any]], str]:
    """Match a proposed column name to a schema column.

    Tries exact -> case-insensitive -> stripped/normalized -> difflib fuzzy.
    Returns (schema_col_dict | None, match_kind).
    """
    if not proposed:
        return None, "none"

    def _norm(s: str) -> str:
        return "".join(ch for ch in s.lower().strip() if ch.isalnum())

    # 1. exact
    for sc in schema_cols:
        if sc["column"] == proposed:
            return sc, "exact"
    # 2. case-insensitive
    pl = proposed.lower().strip()
    for sc in schema_cols:
        if sc["column"].lower().strip() == pl:
            return sc, "ci"
    # 3. normalized (strip spaces/underscores/punct)
    pn = _norm(proposed)
    if pn:
        for sc in schema_cols:
            if _norm(sc["column"]) == pn:
                return sc, "strip"
    # 4. difflib fuzzy
    names = [sc["column"] for sc in schema_cols]
    close = difflib.get_close_matches(proposed, names, n=1, cutoff=_FUZZY_CUTOFF)
    if close:
        for sc in schema_cols:
            if sc["column"] == close[0]:
                return sc, "fuzzy"
    return None, "none"


# --------------------------------------------------------------------------- #
# LLM extraction
# --------------------------------------------------------------------------- #
def _build_prompt(digest: str, schema_cols: List[Dict[str, Any]]) -> str:
    schema_lines = "\n".join(
        f"- {c['table_name']}.{c['column']}" + (f" ({c['dtype']})" if c.get("dtype") else "")
        for c in schema_cols[:200]
    ) or "(no schema columns available)"
    return (
        "You are a data-catalog assistant. You are given (A) the LIVE SCHEMA of a "
        "data source and (B) the text of business definition/explanation documents "
        "(a definitions spreadsheet and/or an explanation deck).\n\n"
        "Extract a STRICT JSON object with EXACTLY these keys:\n"
        '  "column_descriptions": [{"column": "<name as it appears in the docs>", '
        '"description": "<plain-English meaning>"}],\n'
        '  "instructions": [{"content": "<always-on rule / KPI formula / business '
        'rule>", "category": "<kpi|business_rule|compliance|general>"}],\n'
        '  "examples": [{"question": "<natural-language question>", "answer": '
        '"<short answer>", "sql": "<SQL if derivable, else empty string>"}],\n'
        '  "compliance": [{"content": "<compliance / governance rule>"}]\n\n'
        "Rules: Prefer column names that resemble the live schema below. Only emit "
        "facts grounded in the documents. If a section has nothing, return an empty "
        "array for it. Output ONLY the JSON object, no prose, no markdown fences.\n\n"
        "=== LIVE SCHEMA ===\n"
        f"{schema_lines}\n\n"
        "=== DOCUMENTS ===\n"
        f"{digest[:_MAX_DIGEST_CHARS]}\n"
    )


def _coerce_proposal(raw: str) -> Optional[Dict[str, Any]]:
    """Parse the model's reply into the proposal dict; None on junk."""
    if not raw:
        return None
    text = raw.strip()
    # Strip markdown fences if present.
    if text.startswith("```"):
        text = text.strip("`")
        nl = text.find("\n")
        if nl != -1:
            text = text[nl + 1:]
    # Trim to the outermost JSON object.
    start = text.find("{")
    end = text.rfind("}")
    if start != -1 and end != -1 and end > start:
        text = text[start:end + 1]
    try:
        obj = json.loads(text, strict=False)
    except Exception:
        return None
    if not isinstance(obj, dict):
        return None
    return obj


async def _llm_extract(
    db: AsyncSession,
    organization,
    digest: str,
    schema_cols: List[Dict[str, Any]],
) -> Optional[Dict[str, Any]]:
    """One cheap OpenRouter call. Returns the raw proposal dict or None.

    Mirrors the ambiguity-gate / knowledge-proposer idiom exactly.
    """
    try:
        from app.services.llm_service import LLMService

        model = await LLMService().get_default_model(db, organization, None, is_small=True)
        if model is None:
            return None

        from app.ai.llm.llm import LLM
        from app.dependencies import async_session_maker

        prompt = _build_prompt(digest, schema_cols)
        raw = (
            LLM(model, usage_session_maker=async_session_maker).inference(
                prompt, usage_scope="auto_configure_doc"
            )
            or ""
        )
        return _coerce_proposal(raw)
    except Exception as e:
        logger.warning("doc_extractor LLM extract failed: %s", e)
        return None


# --------------------------------------------------------------------------- #
# Offline deterministic fallback
# --------------------------------------------------------------------------- #
def _fallback_from_pairs(pairs: List[Tuple[str, str]]) -> Dict[str, Any]:
    """Build a proposal from a 2-col definitions sheet without an LLM.

    Drops an obvious header row (e.g. 'Column' | 'Definition').
    """
    cds: List[Dict[str, str]] = []
    seen = set()
    for left, right in pairs:
        l, r = (left or "").strip(), (right or "").strip()
        if not l or not r:
            continue
        if l.lower() in ("column", "field", "name", "column name") and \
                r.lower() in ("definition", "description", "meaning", "desc"):
            continue  # header
        key = l.lower()
        if key in seen:
            continue
        seen.add(key)
        cds.append({"column": l, "description": r})
    return {
        "column_descriptions": cds,
        "instructions": [],
        "examples": [],
        "compliance": [],
    }


# --------------------------------------------------------------------------- #
# Public entry point
# --------------------------------------------------------------------------- #
async def extract_proposal(
    db: AsyncSession,
    *,
    organization,
    file_ids: List[str],
    data_source_id: str,
) -> Dict[str, Any]:
    """Parse the docs, extract a proposal, fuzzy-match to the live schema.

    Returns the proposal dict described in the module docstring, or a dict with
    an ``error`` key (route returns 422) for content problems.
    """
    warnings: List[str] = []

    # 1. Validate the target data source (org-scoped).
    ds_res = await db.execute(
        select(DataSource).where(
            DataSource.id == str(data_source_id),
            DataSource.organization_id == organization.id,
        )
    )
    if ds_res.scalar_one_or_none() is None:
        return {"error": "data source not found in this organization"}

    if not file_ids:
        return {"error": "no file_ids provided"}

    # 2. Resolve + parse each file, merging digests/pairs.
    digests: List[str] = []
    pairs: List[Tuple[str, str]] = []
    for fid in file_ids:
        f_res = await db.execute(
            select(FileModel).where(
                FileModel.id == str(fid),
                FileModel.organization_id == organization.id,
            )
        )
        file = f_res.scalar_one_or_none()
        if file is None:
            warnings.append(f"file {fid} not found in org")
            continue
        abs_path = _resolve_upload_path(file.path or "")
        if not abs_path:
            warnings.append(f"file {fid} content missing on disk")
            continue
        parsed = _parse_file(abs_path)
        if "error" in parsed:
            # A missing parser lib is a hard, helpful error for the whole call.
            if "not installed" in parsed["error"]:
                return {"error": parsed["error"]}
            warnings.append(f"file {fid}: {parsed['error']}")
            continue
        if parsed.get("digest"):
            digests.append(parsed["digest"])
        if parsed.get("pairs"):
            pairs.extend(parsed["pairs"])

    digest = "\n\n".join(digests).strip()
    if not digest and not pairs:
        return {"error": "no readable content found in the uploaded files",
                "warnings": warnings}

    # 3. Load the live schema columns to match against.
    schema_cols = await _load_schema_columns(db, data_source_id)

    # 4. LLM extraction, with offline fallback.
    proposal = await _llm_extract(db, organization, digest, schema_cols)
    source = "llm"
    if proposal is None:
        proposal = _fallback_from_pairs(pairs)
        source = "fallback"
        if not pairs:
            warnings.append(
                "no LLM available and no 2-column definitions sheet to parse"
            )

    # 5. Normalize the proposal shape defensively.
    raw_cds = proposal.get("column_descriptions") or []
    raw_instr = proposal.get("instructions") or []
    raw_examples = proposal.get("examples") or []
    raw_compliance = proposal.get("compliance") or []

    # 6. Fuzzy-match each proposed column to the live schema.
    column_descriptions: List[Dict[str, Any]] = []
    unmatched: List[str] = []
    for item in raw_cds:
        if not isinstance(item, dict):
            continue
        col = (item.get("column") or "").strip()
        desc = (item.get("description") or "").strip()
        if not col:
            continue
        match, kind = _match_column(col, schema_cols)
        entry = {
            "column": col,
            "description": desc,
            "matched": match is not None,
            "match_kind": kind,
            "table_id": match["table_id"] if match else None,
            "table_name": match["table_name"] if match else None,
            "matched_column": match["column"] if match else None,
        }
        column_descriptions.append(entry)
        if match is None:
            unmatched.append(col)

    instructions = [
        {"content": (i.get("content") or "").strip(),
         "category": (i.get("category") or "general").strip()}
        for i in raw_instr if isinstance(i, dict) and (i.get("content") or "").strip()
    ]
    examples = [
        {"question": (e.get("question") or "").strip(),
         "answer": (e.get("answer") or "").strip(),
         "sql": (e.get("sql") or "").strip()}
        for e in raw_examples
        if isinstance(e, dict) and (e.get("question") or "").strip()
    ]
    compliance = [
        {"content": (c.get("content") or "").strip()}
        for c in raw_compliance if isinstance(c, dict) and (c.get("content") or "").strip()
    ]

    return {
        "column_descriptions": column_descriptions,
        "instructions": instructions,
        "examples": examples,
        "compliance": compliance,
        "unmatched_columns": unmatched,
        "schema_columns": [f"{c['table_name']}.{c['column']}" for c in schema_cols],
        "source": source,
        "warnings": warnings,
    }


# --------------------------------------------------------------------------- #
# T3 — bind a sibling "Definitions" source onto a data source's column meanings
# --------------------------------------------------------------------------- #
# A "Definitions" upload is a column glossary (column-name -> meaning) that today
# lands as just another queryable table. This wires it so its definitions get
# APPLIED to the *sibling* real data source's column meanings (pending, gated).
#
# Detection of a definitions doc (either signal is enough):
#   * NAME — the data source name or the file name contains a glossary token
#     ('definition', 'glossary', 'dictionary', 'data dict', 'codebook').
#   * SHAPE — a 2-column name->meaning sheet (most rows are exactly 2 cells, the
#     right cell reads like prose), detected by re-using ``_parse_file``.
#
# Bind flow: parse the definitions file(s) via ``extract_proposal`` (which already
# fuzzy-matches proposed column names to the TARGET source's live schema with the
# difflib cutoff above) -> write each MATCHED meaning to the SemanticColumn path
# as ``status='pending'`` (approval-gated, NEVER overwrites an approved row).

_DEFINITION_NAME_TOKENS = (
    "definition", "definitions", "glossary", "dictionary",
    "data dict", "data-dict", "datadict", "codebook", "data catalog",
)


def _name_is_definitions(*names: Optional[str]) -> bool:
    pool = " ".join((n or "").lower() for n in names)
    return any(tok in pool for tok in _DEFINITION_NAME_TOKENS)


def _shape_is_definitions(parsed: Dict[str, Any]) -> bool:
    """Heuristic: does a parsed workbook look like a 2-col name->meaning glossary?

    Conservative — only used as a *fallback* signal when the name doesn't match.
    Pending + approval-gated downstream, so a rare false positive is low-harm.
    """
    pairs = parsed.get("pairs") or []
    if len(pairs) < 3:
        return False
    left_lens, right_lens, prose = [], [], 0
    for left, right in pairs:
        l, r = (left or "").strip(), (right or "").strip()
        if not l or not r:
            continue
        left_lens.append(len(l))
        right_lens.append(len(r))
        if len(r) >= 12 and " " in r:  # meaning cell reads like prose
            prose += 1
    if len(right_lens) < 3:
        return False
    avg_left = sum(left_lens) / len(left_lens)
    avg_right = sum(right_lens) / len(right_lens)
    # Short left (column-like) + longer prose-y right + a majority prose rows.
    return avg_left <= 40 and avg_right >= 12 and prose >= max(3, len(right_lens) // 2)


async def _find_definition_file_ids(
    db: AsyncSession, *, organization, exclude_data_source_id: str
) -> List[str]:
    """Find sibling definitions-doc file ids in this org (excludes the target ds).

    Scans the org's data sources (and their attached files) plus standalone org
    files; classifies each candidate by NAME token or, for .xlsx, 2-col SHAPE.
    Bounded + fail-soft — returns [] on any trouble.
    """
    out: List[str] = []
    try:
        # Candidate map: file_id -> {path, filename, ds_names}
        cands: Dict[str, Dict[str, Any]] = {}

        ds_rows = (
            await db.execute(
                select(DataSource).where(
                    DataSource.organization_id == organization.id
                )
            )
        ).scalars().all()
        for ds in ds_rows:
            if str(ds.id) == str(exclude_data_source_id):
                continue
            for f in (getattr(ds, "files", None) or []):
                ent = cands.setdefault(
                    str(f.id),
                    {"path": f.path, "filename": f.filename, "ds_names": set()},
                )
                ent["ds_names"].add(ds.name or "")

        # Standalone org files (covers a definitions file not wrapped as a ds).
        f_rows = (
            await db.execute(
                select(FileModel).where(
                    FileModel.organization_id == organization.id
                )
            )
        ).scalars().all()
        for f in f_rows:
            cands.setdefault(
                str(f.id),
                {"path": f.path, "filename": f.filename, "ds_names": set()},
            )

        parsed_budget = 25  # cap content-shape parses (name hits are free)
        for fid, ent in cands.items():
            names = list(ent["ds_names"]) + [ent.get("filename")]
            if _name_is_definitions(*names):
                out.append(fid)
                continue
            # SHAPE fallback — only for spreadsheets, only while budget remains.
            fname = (ent.get("filename") or "").lower()
            if parsed_budget <= 0 or not fname.endswith((".xlsx", ".xlsm")):
                continue
            abs_path = _resolve_upload_path(ent.get("path") or "")
            if not abs_path:
                continue
            parsed_budget -= 1
            parsed = _parse_file(abs_path)
            if "error" not in parsed and _shape_is_definitions(parsed):
                out.append(fid)
    except Exception as e:  # never raise — this feeds a fail-soft train stage
        logger.warning("doc_extractor _find_definition_file_ids failed: %s", e)
    return out


async def apply_definitions_to_data_source(
    db: AsyncSession, *, organization, data_source, model=None
) -> Dict[str, Any]:
    """Bind a sibling Definitions doc onto a data source's column meanings.

    Finds a definitions-glossary upload in the same org, parses + fuzzy-matches
    its terms to THIS data source's columns (via ``extract_proposal``), and writes
    each matched meaning to the existing ``SemanticColumn`` path as
    ``status='pending'`` (approval-gated; NEVER overwrites an approved row; seeds
    the SemanticTable/SemanticColumn skeleton if it isn't there yet).

    Gated behind the EXISTING ``SEMANTIC_LAYER`` flag. NEVER raises — degrades to
    ``{'columns': [], ...}`` so it can never break Auto-train.
    Returns ``{'columns': [<SemanticColumn.id>...], 'matched': int, 'source': ...}``.
    """
    out: Dict[str, Any] = {"columns": [], "matched": 0}
    try:
        from app.settings.hybrid_flags import flags
        if not flags.SEMANTIC_LAYER:
            return out

        ds_id = str(getattr(data_source, "id", None) or "")
        if not ds_id:
            return out

        file_ids = await _find_definition_file_ids(
            db, organization=organization, exclude_data_source_id=ds_id
        )
        if not file_ids:
            return out

        # Parse + fuzzy-match the definitions docs to THIS ds's live schema.
        proposal = await extract_proposal(
            db,
            organization=organization,
            file_ids=file_ids,
            data_source_id=ds_id,
        )
        if not isinstance(proposal, dict) or "error" in proposal:
            out["error"] = (proposal or {}).get("error") if isinstance(proposal, dict) else "no proposal"
            return out
        out["source"] = proposal.get("source")

        matched = [
            cd for cd in (proposal.get("column_descriptions") or [])
            if isinstance(cd, dict) and cd.get("matched")
            and (cd.get("description") or "").strip()
            and cd.get("table_name") and cd.get("matched_column")
        ]
        out["matched"] = len(matched)
        if not matched:
            return out

        from app.models.semantic_table import SemanticTable, SemanticColumn
        from sqlalchemy.orm import selectinload

        # Load existing semantic skeleton for (org, ds) — columns eager-loaded.
        sem_tables = (
            await db.execute(
                select(SemanticTable)
                .where(
                    SemanticTable.organization_id == str(organization.id),
                    SemanticTable.data_source_id == ds_id,
                )
                .options(selectinload(SemanticTable.columns))
            )
        ).scalars().all()
        st_by_name: Dict[str, Any] = {st.table_name: st for st in sem_tables}
        cols_by_table: Dict[str, Dict[str, Any]] = {
            st.table_name: {c.name: c for c in (st.columns or [])}
            for st in sem_tables
        }

        changed = False
        for cd in matched:
            tname = cd["table_name"]
            cname = cd["matched_column"]
            meaning = (cd["description"] or "").strip()

            st = st_by_name.get(tname)
            if st is None:  # seed the table skeleton (mirrors GET /semantic)
                st = SemanticTable(
                    organization_id=str(organization.id),
                    data_source_id=ds_id,
                    table_name=tname,
                    description="",
                    use_cases=[],
                    quality_notes=[],
                    status="draft",
                )
                db.add(st)
                await db.flush()
                st_by_name[tname] = st
                cols_by_table[tname] = {}

            col = cols_by_table.get(tname, {}).get(cname)
            if col is None:  # seed the column skeleton
                col = SemanticColumn(
                    semantic_table_id=st.id,
                    name=cname,
                    type="",
                    meaning="",
                    status="draft",
                )
                db.add(col)
                await db.flush()
                cols_by_table.setdefault(tname, {})[cname] = col

            # NEVER overwrite an approved meaning; skip no-op rewrites.
            if _clean(getattr(col, "status", None)) == "approved":
                continue
            if _clean(getattr(col, "meaning", None)) == meaning:
                continue
            col.meaning = meaning
            col.status = "pending"
            out["columns"].append(str(col.id))
            changed = True

        if changed:
            await db.commit()
        return out
    except Exception as e:  # never raise to the caller (train/route is fail-soft)
        logger.warning("doc_extractor apply_definitions_to_data_source failed: %s", e)
        try:
            await db.rollback()
        except Exception:
            pass
        return {"columns": [], "matched": 0, "error": str(e)}


def _clean(v: Any) -> str:
    return str(v or "").strip()
